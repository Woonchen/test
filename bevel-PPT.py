#The Auto_report_v1.0 for bevel trimming by Nick Huang.

from pptx import Presentation
from pptx.util import Pt, Cm 
from pptx.dml.color import RGBColor
import datetime



prs = Presentation('/Users/nick/Desktop/nick/code/python/sample.pptx')

img_path = '/Users/nick/Desktop/project/AA_晶技/2AAA01-bevel trimming/2AAA01T007-10片/加工前/1-BB.bmp'
img_path1 = '/Users/nick/Desktop/nick/code/python/position.png'
Today = datetime.date.today()


#新增投影片[母片序號]
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.add_textbox(left=Cm(4.86), top=Cm(9.38), width=Cm(24.35), height=Cm(1.45)).text_frame
title.paragraphs[0].text = "2AAA01-TXC Quartz Wafer Bevel Profile Trimming Project"
title.paragraphs[0].font.size = Pt(28)
title.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

sub = slide.shapes.add_textbox(left=Cm(4.86), top=Cm(11.52), width=Cm(16.92), height=Cm(1.8)).text_frame
sub.paragraphs[0].text = "專案負責人：Nick\n更新日期："+ str(Today)
sub.paragraphs[0].font.size = Pt(18)
sub.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)


slide = prs.slides.add_slide(prs.slide_layouts[2])
#插入圖片
pic_1 = slide.shapes.add_picture(img_path, left=Cm(2.7), top=Cm(3.06), width=Cm(9.18), height=Cm(6.8)) #在slide這頁插入圖片
pic_1.crop_top = 1.38/9.18
pic_1.crop_bottom = 1/9.18
pic_2 = slide.shapes.add_picture(img_path, left=Cm(2.7), top=Cm(10.94), width=Cm(9.18), height=Cm(6.8)) #在slide這頁插入圖片
pic_2.crop_top = 1.38/9.18
pic_2.crop_bottom = 1/9.18
pic_3 = slide.shapes.add_picture(img_path, left=Cm(14.7), top=Cm(2.8), width=Cm(14.69), height=Cm(6.8)) #在slide這頁插入圖片
pic_4 = slide.shapes.add_picture(img_path, left=Cm(14.7), top=Cm(13), width=Cm(10.47), height=Cm(3.15)) #在slide這頁插入圖片
pic_5 = slide.shapes.add_picture(img_path1, left=Cm(27.12), top=Cm(12.19), width=Cm(4.55), height=Cm(5.15)) #在slide這頁插入圖片

title = slide.shapes.add_textbox(left=Cm(1.35), top=Cm(0.45), width=Cm(32.1), height=Cm(1.23)).text_frame
title.paragraphs[0].text = "Measurement Date  #1 (Position 1)"
title.paragraphs[0].font.size = Pt(28)
str_before = slide.shapes.add_textbox(left=Cm(5.11), top=Cm(2), width=Cm(4.77), height=Cm(1.03))
str_before.text = "Before trimming"
str_after = slide.shapes.add_textbox(left=Cm(5.11), top=Cm(9.86), width=Cm(4.77), height=Cm(1.03))
str_after.text = "After trimming"


prs.save('test'+str(Today)+'.pptx')








# XX= str(input("(20200710):"))