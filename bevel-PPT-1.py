#The Auto_report_v1.1 for bevel trimming by Nick Huang.

from pptx import Presentation
from pptx.util import Pt, Cm 
from pptx.dml.color import RGBColor
import datetime

caseNumber = "2AAA01T014"#str(input("caseNumber:"))
casePath = "/Users/nick/Desktop/project/AA_晶技/2AAA01-bevel trimming/2AAA01T014/圖"#str(input("casePath:"))

prs = Presentation('/Users/nick/Desktop/nick/code/python/ppt/sample.pptx')
sample_path1 = '/Users/nick/Desktop/nick/code/python/ppt/position.png'
Today = datetime.date.today()

slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.add_textbox(left=Cm(4.86), top=Cm(9.38), width=Cm(24.35), height=Cm(1.45)).text_frame
title.paragraphs[0].text = caseNumber+"-TXC Quartz Wafer Bevel Profile Trimming Project"
title.paragraphs[0].font.size = Pt(28)
title.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
sub = slide.shapes.add_textbox(left=Cm(4.86), top=Cm(11.52), width=Cm(16.92), height=Cm(1.8)).text_frame
sub.paragraphs[0].text = "專案負責人：Nick\n更新日期："+ str(Today)
sub.paragraphs[0].font.size = Pt(18)
sub.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
slide = prs.slides.add_slide(prs.slide_layouts[1])


def page(number):
    for position in ["L", "T", "R"]:   # "B"
        imgPath_B = casePath +'/'+ str(number) +'/'+str(number) +'-B'+ position +'.bmp'
        imgPath_A = casePath +'/'+ str(number) +'/'+str(number) +'-A'+ position +'.bmp'
        imgPath_XS = '/Users/nick/Desktop/project/AA_晶技/2AAA01-bevel trimming/2AAA01T014/圖/換算圖'+'/'+str(number) +'-'+ position +'.png'
        imgPath_table = '/Users/nick/Desktop/project/AA_晶技/2AAA01-bevel trimming/2AAA01T014/圖/換算圖'+'/'+str(number) +'-'+ position +'-表格.png'

        if position == "B":
            positionNo = "1"
        elif position == "L":
            positionNo = "2"
        elif position == "T":
            positionNo = "3"
        elif position == "R":
            positionNo = "4"

        slide = prs.slides.add_slide(prs.slide_layouts[2])
        imgBefore = slide.shapes.add_picture(imgPath_B, left=Cm(2.7), top=Cm(3.06), width=Cm(9.18), height=Cm(6.8)) 
        imgBefore.crop_top = 1.38/9.18
        imgBefore.crop_bottom = 1/9.18
        imgAfter = slide.shapes.add_picture(imgPath_A, left=Cm(2.7), top=Cm(10.94), width=Cm(9.18), height=Cm(6.8)) 
        imgAfter.crop_top = 1.38/9.18
        imgAfter.crop_bottom = 1/9.18
        XS = slide.shapes.add_picture(imgPath_XS, left=Cm(14.7), top=Cm(2.8), width=Cm(14.69), height=Cm(6.8)) 
        tabel = slide.shapes.add_picture(imgPath_table, left=Cm(14.7), top=Cm(13), width=Cm(10.47), height=Cm(3.15)) 
        sample = slide.shapes.add_picture(sample_path1, left=Cm(27.12), top=Cm(12.19), width=Cm(4.55), height=Cm(5.15)) 
        title = slide.shapes.add_textbox(left=Cm(1.35), top=Cm(0.45), width=Cm(32.1), height=Cm(1.23)).text_frame
        title.paragraphs[0].text = "Measurement Date  #"+ str(number) + " (Position " + positionNo + ")"
        title.paragraphs[0].font.size = Pt(28)
        str_before = slide.shapes.add_textbox(left=Cm(5.11), top=Cm(2), width=Cm(4.77), height=Cm(1.03))
        str_before.text = "Before trimming"
        str_after = slide.shapes.add_textbox(left=Cm(5.11), top=Cm(9.86), width=Cm(4.77), height=Cm(1.03))
        str_after.text = "After trimming"

#

page(1)
page(2)
page(3)
page(4)
page(5)
page(6)
page(7)
page(8)
page(9)
page(10)
page("研磨")
page("研磨-4")
page("研磨-5")





prs.save('/Users/nick/Desktop/'+caseNumber+'.pptx')


# /Users/nick/Desktop/bevel/2AAA01T010

# XX= str(input("(20200710):"))

#/Users/nick/Desktop/project/AA_晶技/2AAA01-bevel trimming/2AAA01T011/石英邊緣量測-230105/Users/nick/Desktop/project/AA_晶技/2AAA01-bevel trimming/2AAA01T011/石英邊緣量測-230105